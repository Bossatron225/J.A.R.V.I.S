"""
file_processor.py — JARVIS Universal File Processor

Supported types:
  image   → describe, ocr, resize, convert, compress, crop
  pdf     → summarize, extract_text, extract_pages, to_word
  docx    → summarize, extract_text, reformat, translate_hint
  txt/md  → summarize, reformat, translate_hint, word_count
  csv     → analyze, filter, sort, convert, stats
  xlsx    → analyze, filter, convert, stats
  json    → validate, format, extract, convert
  code    → explain, review, fix, run, document
  audio   → transcribe, trim, convert, info
  video   → trim, extract_audio, extract_frame, info, compress
  zip     → list, extract
  pptx    → summarize, extract_text, to_pdf
"""

import os
import io
import re
import json
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from datetime import datetime

def _get_api_key() -> str:
    config_path = Path(__file__).resolve().parent.parent / "config" / "api_keys.json"
    with open(config_path, "r", encoding="utf-8") as f:
        return json.load(f)["gemini_api_key"]


def _gemini_client():
    from google import genai
    _c = genai.Client(api_key=_get_api_key())

    class _W:
        def generate_content(self, contents):
            return _c.models.generate_content(model="models/gemini-flash-lite-latest", contents=contents)

    return _W()


def _detect_type(path: Path) -> str:
    ext = path.suffix.lower().lstrip(".")
    image_exts = {"jpg", "jpeg", "png", "gif", "webp", "bmp", "tiff", "svg", "ico"}
    video_exts = {"mp4", "avi", "mov", "mkv", "wmv", "flv", "webm", "m4v", "3gp"}
    audio_exts = {"mp3", "wav", "ogg", "m4a", "aac", "flac", "wma", "opus"}
    code_exts  = {"py", "js", "ts", "jsx", "tsx", "html", "css", "java", "c",
                  "cpp", "cs", "go", "rs", "rb", "php", "swift", "kt", "sh",
                  "bash", "ps1", "lua", "r", "m", "sql", "yaml", "toml"}
    archive_exts = {"zip", "rar", "tar", "gz", "7z", "bz2", "xz"}

    if ext in image_exts:  return "image"
    if ext in video_exts:  return "video"
    if ext in audio_exts:  return "audio"
    if ext in code_exts:   return "code"
    if ext in archive_exts: return "archive"
    if ext == "pdf":       return "pdf"
    if ext in ("docx", "doc"): return "docx"
    if ext in ("txt", "md", "rst", "log"): return "text"
    if ext in ("csv", "tsv"): return "csv"
    if ext in ("xlsx", "xls", "ods"): return "excel"
    if ext == "json":      return "json"
    if ext == "xml":       return "xml"
    if ext in ("pptx", "ppt"): return "pptx"
    return "unknown"


def _file_size_str(path: Path) -> str:
    size = path.stat().st_size
    if size < 1024:        return f"{size} B"
    if size < 1024**2:     return f"{size/1024:.1f} KB"
    if size < 1024**3:     return f"{size/1024**2:.1f} MB"
    return f"{size/1024**3:.1f} GB"

def _output_path(src: Path, suffix: str, new_ext: str = None) -> Path:
    ext  = new_ext or src.suffix
    name = f"{src.stem}_{suffix}{ext}"
    return src.parent / name


def _image_prompt_for_context(context: str, path: Path, name: str | None = None) -> str:
    label = name or path.name
    return (
        f"Analyse the visual content in this {context} image from {label}. "
        "Describe the subject, any visible text, charts, UI elements, diagrams, labels, and notable details. "
        "Keep the answer concise but information dense."
    )


def _summarize_pil_image(image, prompt: str, instruction: str | None = None) -> str:
    try:
        model = _gemini_client()
        final_prompt = instruction.strip() if instruction and instruction.strip() else prompt
        response = model.generate_content([final_prompt, image])
        return response.text.strip()
    except Exception as e:
        return f"Image analysis failed: {e}"


def _load_pil_image(data: bytes):
    from PIL import Image

    image = Image.open(io.BytesIO(data))
    if image.mode not in ("RGB", "RGBA"):
        image = image.convert("RGB")
    return image


def _extract_docx_text_and_images(path: Path) -> tuple[str, list[str]]:
    text = ""
    try:
        from docx import Document

        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        table_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_text:
                    table_text.append(" | ".join(row_text))
        text = "\n".join(paragraphs + table_text)
    except Exception:
        text = path.read_text(encoding="utf-8", errors="ignore")

    image_notes: list[str] = []
    try:
        with zipfile.ZipFile(path) as archive:
            media_files = [name for name in archive.namelist() if name.startswith("word/media/") and not name.endswith("/")]
            for index, media_name in enumerate(media_files[:8], 1):
                try:
                    image = _load_pil_image(archive.read(media_name))
                    note = _summarize_pil_image(
                        image,
                        _image_prompt_for_context("DOCX", path, f"embedded image {index}"),
                    )
                    if note:
                        image_notes.append(f"[Image {index}: {Path(media_name).name}]\n{note}")
                except Exception as e:
                    image_notes.append(f"[Image {index}: {Path(media_name).name}]\nImage analysis failed: {e}")
    except Exception:
        pass

    return text, image_notes


def _extract_pptx_text_and_images(path: Path) -> tuple[str, list[str]]:
    text = ""
    try:
        from pptx import Presentation

        prs = Presentation(path)
        slide_text = []
        for i, slide in enumerate(prs.slides, 1):
            block = [f"--- Slide {i} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    block.append(shape.text.strip())
            slide_text.append("\n".join(block))
        text = "\n\n".join(slide_text)
    except Exception:
        text = ""

    image_notes: list[str] = []
    try:
        with zipfile.ZipFile(path) as archive:
            media_files = [name for name in archive.namelist() if name.startswith("ppt/media/") and not name.endswith("/")]
            for index, media_name in enumerate(media_files[:8], 1):
                try:
                    image = _load_pil_image(archive.read(media_name))
                    note = _summarize_pil_image(
                        image,
                        _image_prompt_for_context("PPTX", path, f"embedded image {index}"),
                    )
                    if note:
                        image_notes.append(f"[Image {index}: {Path(media_name).name}]\n{note}")
                except Exception as e:
                    image_notes.append(f"[Image {index}: {Path(media_name).name}]\nImage analysis failed: {e}")
    except Exception:
        pass

    return text, image_notes


def _render_pdf_pages_as_images(path: Path, page_limit: int = 8) -> list[object]:
    try:
        import pymupdf
    except Exception:
        return []

    images = []
    try:
        doc = pymupdf.open(path)
        limit = min(len(doc), page_limit)
        for page_index in range(limit):
            page = doc[page_index]
            pix = page.get_pixmap(matrix=pymupdf.Matrix(1.7, 1.7), alpha=False)
            image = _load_pil_image(pix.tobytes("png"))
            images.append(image)
        doc.close()
    except Exception:
        return []
    return images


def _extract_pdf_text_and_visual_notes(path: Path, page_limit: int = 8) -> tuple[str, list[str]]:
    text = ""
    try:
        import pymupdf

        doc = pymupdf.open(path)
        text_parts = []
        for page in doc:
            page_text = page.get_text("text") or ""
            if page_text.strip():
                text_parts.append(page_text.strip())
        text = "\n\n".join(text_parts)
        doc.close()
    except Exception:
        text = ""

    image_notes: list[str] = []
    page_images = _render_pdf_pages_as_images(path, page_limit=page_limit if not text.strip() else max(3, min(page_limit, 3)))
    for index, image in enumerate(page_images, 1):
        note = _summarize_pil_image(
            image,
            _image_prompt_for_context("PDF page", path, f"page {index}"),
        )
        if note:
            image_notes.append(f"[Page {index}]\n{note}")

    return text, image_notes


def _merge_text_and_image_notes(text: str, image_notes: list[str]) -> str:
    parts = []
    if text.strip():
        parts.append(text.strip())
    if image_notes:
        parts.append("Image-based content:\n" + "\n\n".join(image_notes))
    return "\n\n".join(parts).strip()


def _essay_style_brief(style: str, topic: str) -> str:
    style_key = (style or "student").strip().lower()
    styles = {
        "student": (
            f"Write a polished student essay about this {topic}. "
            "Use clear prose, not bullet points. Include a short title, an opening thesis, "
            "2 to 4 body paragraphs, and a brief concluding paragraph. "
            "Keep the tone natural, organized, and easy to follow."
        ),
        "professional": (
            f"Write a polished professional report-style essay about this {topic}. "
            "Use clear prose, not bullet points. Include a short title, an opening thesis, "
            "2 to 4 body paragraphs, and a brief concluding paragraph. "
            "Keep the tone concise, analytical, and formal."
        ),
        "casual": (
            f"Write a polished casual explanation about this {topic}. "
            "Use clear prose, not bullet points. Include a short title, a short opening, "
            "2 to 4 short body paragraphs, and a brief closing paragraph. "
            "Keep the tone friendly, readable, and straightforward."
        ),
    }
    return styles.get(style_key, styles["student"])


def _essay_prompt(topic: str, content: str, instruction: str = "", style: str = "student") -> str:
    base = _essay_style_brief(style, topic)
    extra = instruction.strip()
    if extra:
        base = f"{base} {extra}"
    return f"{base}\n\nSource material:\n{content}"

def _process_image(path: Path, action: str, params: dict, speak=None) -> str:
    try:
        from PIL import Image
    except ImportError:
        return "Pillow is not installed. Run: pip install Pillow"

    action = action or "describe"

    if action in ("describe", "ocr", "analyze", "read", "extract_text"):
        try:
            model  = _gemini_client()
            img    = Image.open(path)
            prompt = {
                "describe": "Describe this image in detail.",
                "ocr":      "Extract all text visible in this image. Return only the text, formatted clearly.",
                "analyze":  "Analyze this image thoroughly: objects, colors, composition, any text, context.",
                "read":     "Read all text in this image, preserving structure and formatting.",
                "extract_text": "Extract all text from this image.",
            }.get(action, "Describe this image.")

            if params.get("instruction"):
                prompt = params["instruction"]

            response = model.generate_content([prompt, img])
            result   = response.text.strip()

            if len(result) > 500 and params.get("save", True):
                out = _output_path(path, "result", ".txt")
                out.write_text(result, encoding="utf-8")
                return f"{result}\n\nFull result saved to: {out}"
            return result
        except Exception as e:
            return f"AI image analysis failed: {e}"

    if action == "resize":
        width  = int(params.get("width",  0))
        height = int(params.get("height", 0))
        scale  = float(params.get("scale", 0))
        try:
            img = Image.open(path)
            w, h = img.size
            if scale:
                new_size = (int(w * scale), int(h * scale))
            elif width and height:
                new_size = (width, height)
            elif width:
                new_size = (width, int(h * width / w))
            elif height:
                new_size = (int(w * height / h), height)
            else:
                return "Please specify width, height, or scale."
            out = _output_path(path, f"resized_{new_size[0]}x{new_size[1]}")
            img.resize(new_size, Image.LANCZOS).save(out)
            return f"Resized from {w}x{h} to {new_size[0]}x{new_size[1]}. Saved: {out.name}"
        except Exception as e:
            return f"Resize failed: {e}"

    if action == "convert":
        fmt = params.get("format", "png").lower().strip(".")
        fmt_map = {"jpg": "JPEG", "jpeg": "JPEG", "png": "PNG",
                   "webp": "WEBP", "bmp": "BMP", "tiff": "TIFF"}
        pil_fmt = fmt_map.get(fmt, fmt.upper())
        try:
            img = Image.open(path).convert("RGB") if fmt == "jpg" else Image.open(path)
            out = _output_path(path, "converted", f".{fmt}")
            img.save(out, pil_fmt)
            return f"Converted to {fmt.upper()}. Saved: {out.name}"
        except Exception as e:
            return f"Convert failed: {e}"

    if action == "compress":
        quality = int(params.get("quality", 70))
        try:
            img = Image.open(path).convert("RGB")
            out = _output_path(path, f"compressed_q{quality}", ".jpg")
            img.save(out, "JPEG", quality=quality, optimize=True)
            before = _file_size_str(path)
            after  = _file_size_str(out)
            return f"Compressed: {before} → {after}. Saved: {out.name}"
        except Exception as e:
            return f"Compress failed: {e}"

    if action == "info":
        try:
            img = Image.open(path)
            return (f"Image info: {img.format}, {img.size[0]}x{img.size[1]}px, "
                    f"mode: {img.mode}, size: {_file_size_str(path)}")
        except Exception as e:
            return f"Info failed: {e}"

    return _process_image(path, "describe", {"instruction": f"{action}: {params}"})

def _process_pdf(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "summarize"
    style = params.get("style") or params.get("tone") or params.get("essay_style") or "student"

    if action in ("summarize", "extract_text", "translate_hint", "analyze", "reformat"):
        text, image_notes = _extract_pdf_text_and_visual_notes(path)
        combined = _merge_text_and_image_notes(text[:50000], image_notes)
        if not combined.strip():
            return "Could not extract text or visual content from PDF."

        if action == "extract_text":
            out = _output_path(path, "text", ".txt")
            out.write_text(combined, encoding="utf-8")
            return f"Text and visual notes extracted ({len(combined)} chars). Saved: {out.name}"

        prompt_map = {
            "summarize":      _essay_prompt("PDF document", combined, "Write a polished essay summary of this PDF. Keep it concise but complete. Use an opening thesis, a few body paragraphs, and a short conclusion.", style),
            "analyze":        _essay_prompt("PDF document", combined, "Write a deeper analytical essay about this PDF. Focus on themes, structure, intent, and important supporting details.", style),
            "translate_hint": _essay_prompt("PDF document", combined, "Write a short essay explaining the language, meaning, and overall content of this PDF.", style),
            "reformat":       _essay_prompt("PDF document", combined, "Rewrite the material as a clean essay with clear paragraphs and natural flow.", style),
        }
        try:
            model    = _gemini_client()
            response = model.generate_content(prompt_map.get(action, f"Analyze:\n\n{text}"))
            result   = response.text.strip()
            if len(result) > 600 and params.get("save", True):
                out = _output_path(path, action, ".txt")
                out.write_text(result, encoding="utf-8")
                return f"{result}\n\nFull result saved: {out.name}"
            return result
        except Exception as e:
            return f"AI analysis failed: {e}"

    if action == "info":
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                pages = len(pdf.pages)
            return f"PDF: {pages} pages, size: {_file_size_str(path)}"
        except Exception:
            return f"PDF size: {_file_size_str(path)}"

    if action == "to_word":
        text = _extract_pdf_text()
        if not text:
            return "Could not extract text to convert."
        try:
            from docx import Document
            doc  = Document()
            doc.add_heading(path.stem, 0)
            for para in text.split("\n\n"):
                if para.strip():
                    doc.add_paragraph(para.strip())
            out = _output_path(path, "converted", ".docx")
            doc.save(out)
            return f"Converted to Word document. Saved: {out.name}"
        except ImportError:
            return "python-docx not installed. Run: pip install python-docx"

    return f"Unknown PDF action: '{action}'. Try: summarize, extract_text, info, to_word"

def _process_text_doc(path: Path, file_type: str, action: str,
                       params: dict, speak=None) -> str:
    action = action or "summarize"
    style = params.get("style") or params.get("tone") or params.get("essay_style") or "student"

    if file_type == "docx":
        content, image_notes = _extract_docx_text_and_images(path)
    else:
        content, image_notes = path.read_text(encoding="utf-8", errors="ignore"), []

    combined = _merge_text_and_image_notes(content, image_notes)
    if not combined.strip():
        return "File appears to be empty."

    if action == "word_count":
        words = len(combined.split())
        chars = len(combined)
        lines = combined.count("\n")
        return f"Word count: {words} words, {chars} characters, {lines} lines."

    if action == "extract_text":
        if file_type != "txt":
            out = _output_path(path, "extracted", ".txt")
            out.write_text(combined, encoding="utf-8")
            return f"Text and visual notes extracted. Saved: {out.name}"
        return combined[:2000]

    instruction = params.get("instruction", "")
    prompt_map  = {
        "summarize":  _essay_prompt("document", combined[:40000], "Write a polished sample essay that summarizes the document. Keep the voice clear and readable, with a title, thesis, body paragraphs, and conclusion.", style),
        "analyze":    _essay_prompt("document", combined[:40000], "Write a thoughtful analytical essay about the document. Explain the main ideas, tone, structure, and any recurring themes or contrasts.", style),
        "reformat":   _essay_prompt("document", combined[:40000], "Rewrite the content into a clean essay with proper structure, strong transitions, and natural prose.", style),
        "fix":        _essay_prompt("document", combined[:40000], "Improve the grammar, spelling, and style of this material while keeping it essay-like and polished.", style),
        "translate_hint": _essay_prompt("document", combined[:10000], "Write a short essay explaining the language and meaning of the text.", style),
        "to_bullet":  f"Convert this text into a clear bullet-point summary:\n\n{combined[:40000]}",
        "custom":     _essay_prompt("document", combined[:40000], instruction, style),
    }

    if action not in prompt_map:

        action  = "custom"
        instruction = action

    try:
        model    = _gemini_client()
        response = model.generate_content(prompt_map[action])
        result   = response.text.strip()
        if len(result) > 600 and params.get("save", True):
            out = _output_path(path, action, ".txt")
            out.write_text(result, encoding="utf-8")
            return f"{result}\n\nFull result saved: {out.name}"
        return result
    except Exception as e:
        return f"AI processing failed: {e}"


def _process_data(path: Path, file_type: str, action: str,
                  params: dict, speak=None) -> str:
    try:
        import pandas as pd
    except ImportError:
        return "pandas not installed. Run: pip install pandas openpyxl"

    action = action or "analyze"

    try:
        if file_type == "csv":
            df = pd.read_csv(path, encoding="utf-8", errors="replace")
        else:
            df = pd.read_excel(path)
    except Exception as e:
        return f"Could not read file: {e}"

    if action == "info":
        return (f"Rows: {len(df)}, Columns: {len(df.columns)}\n"
                f"Columns: {', '.join(df.columns.tolist())}\n"
                f"Size: {_file_size_str(path)}")

    if action == "stats":
        try:
            desc = df.describe(include="all").to_string()
            return f"Statistics:\n{desc[:2000]}"
        except Exception as e:
            return f"Stats failed: {e}"

    if action == "analyze":
        preview = df.head(50).to_string()
        prompt  = (f"Analyze this dataset. Columns: {list(df.columns)}\n"
                   f"Rows: {len(df)}\nPreview:\n{preview}\n\n"
                   f"Give insights, patterns, and notable findings.")
        try:
            model    = _gemini_client()
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"AI analysis failed: {e}"

    if action in ("convert", "to_csv", "to_excel", "to_json"):
        fmt = {"to_csv": "csv", "to_excel": "xlsx", "to_json": "json",
               "convert": params.get("format", "csv")}.get(action, "csv")
        try:
            if fmt == "csv":
                out = _output_path(path, "converted", ".csv")
                df.to_csv(out, index=False, encoding="utf-8")
            elif fmt == "xlsx":
                out = _output_path(path, "converted", ".xlsx")
                df.to_excel(out, index=False)
            elif fmt == "json":
                out = _output_path(path, "converted", ".json")
                df.to_json(out, orient="records", force_ascii=False, indent=2)
            return f"Converted to {fmt.upper()}. Saved: {out.name}"
        except Exception as e:
            return f"Convert failed: {e}"

    if action == "filter":
        col       = params.get("column", "")
        value     = params.get("value", "")
        condition = params.get("condition", "equals")
        if not col or col not in df.columns:
            return f"Column '{col}' not found. Available: {', '.join(df.columns)}"
        try:
            if condition == "equals":     filtered = df[df[col] == value]
            elif condition == "contains": filtered = df[df[col].astype(str).str.contains(str(value), case=False)]
            elif condition == "gt":       filtered = df[df[col] > float(value)]
            elif condition == "lt":       filtered = df[df[col] < float(value)]
            else:                         filtered = df[df[col] == value]
            out = _output_path(path, "filtered", ".csv")
            filtered.to_csv(out, index=False)
            return f"Filtered: {len(filtered)} rows match. Saved: {out.name}"
        except Exception as e:
            return f"Filter failed: {e}"

    if action == "sort":
        col = params.get("column", df.columns[0])
        asc = params.get("ascending", True)
        try:
            sorted_df = df.sort_values(col, ascending=asc)
            out = _output_path(path, "sorted", path.suffix)
            sorted_df.to_csv(out, index=False)
            return f"Sorted by '{col}'. Saved: {out.name}"
        except Exception as e:
            return f"Sort failed: {e}"

    preview = df.head(30).to_string()
    try:
        model    = _gemini_client()
        response = model.generate_content(
            f"Task: {action}\nDataset ({len(df)} rows, cols: {list(df.columns)}):\n{preview}"
        )
        return response.text.strip()
    except Exception as e:
        return f"Processing failed: {e}"


def _process_json(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "analyze"
    try:
        content = path.read_text(encoding="utf-8")
        data    = json.loads(content)
    except Exception as e:
        return f"Invalid JSON: {e}"

    if action == "validate":
        return f"Valid JSON. Type: {type(data).__name__}, size: {_file_size_str(path)}"

    if action == "format":
        out = _output_path(path, "formatted", ".json")
        out.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
        return f"Formatted JSON saved: {out.name}"

    if action in ("analyze", "summarize", "extract"):
        preview = json.dumps(data, indent=2, ensure_ascii=False)[:8000]
        prompt  = f"Task: {action} this JSON data:\n{preview}"
        if params.get("instruction"):
            prompt = f"{params['instruction']}\n\nJSON data:\n{preview}"
        try:
            model    = _gemini_client()
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"AI processing failed: {e}"

    if action == "to_csv":
        try:
            import pandas as pd
            if isinstance(data, list):
                df  = pd.DataFrame(data)
                out = _output_path(path, "converted", ".csv")
                df.to_csv(out, index=False)
                return f"Converted to CSV. Saved: {out.name}"
            return "JSON must be an array of objects to convert to CSV."
        except ImportError:
            return "pandas not installed."

    return _process_json(path, "analyze", {"instruction": action})

def _process_code(path: Path, action: str, params: dict, speak=None) -> str:
    action  = action or "explain"
    content = path.read_text(encoding="utf-8", errors="ignore")
    ext     = path.suffix.lstrip(".")

    if action == "run":
        if ext == "py":
            try:
                result = subprocess.run(
                    ["python", str(path)],
                    capture_output=True, text=True, timeout=30
                )
                out = result.stdout or result.stderr
                return f"Output:\n{out[:2000]}" if out else "No output."
            except subprocess.TimeoutExpired:
                return "Execution timed out (30s)."
            except Exception as e:
                return f"Run failed: {e}"
        return f"Direct execution not supported for .{ext} files."

    if action == "info":
        lines = content.count("\n")
        words = len(content.split())
        return f"Code file: {lines} lines, {words} words, {_file_size_str(path)}"

    prompt_map = {
        "explain":   f"Explain this {ext} code clearly:\n\n```{ext}\n{content[:30000]}\n```",
        "review":    f"Review this {ext} code for bugs, issues, and improvements:\n\n```{ext}\n{content[:30000]}\n```",
        "fix":       f"Fix any bugs in this {ext} code and return the corrected version:\n\n```{ext}\n{content[:30000]}\n```",
        "optimize":  f"Optimize this {ext} code for performance and readability:\n\n```{ext}\n{content[:30000]}\n```",
        "document":  f"Add proper documentation/comments to this {ext} code:\n\n```{ext}\n{content[:30000]}\n```",
        "summarize": f"Summarize what this {ext} code does:\n\n```{ext}\n{content[:30000]}\n```",
        "test":      f"Write unit tests for this {ext} code:\n\n```{ext}\n{content[:30000]}\n```",
    }

    instruction = params.get("instruction", "")
    if action not in prompt_map:
        prompt = f"{action}\n\n```{ext}\n{content[:30000]}\n```"
        if instruction:
            prompt = f"{instruction}\n\n```{ext}\n{content[:30000]}\n```"
    else:
        prompt = prompt_map[action]

    try:
        model    = _gemini_client()
        response = model.generate_content(prompt)
        result   = response.text.strip()

        if action in ("fix", "optimize", "document") and params.get("save", True):
            out = _output_path(path, action)
            code_match = re.search(r"```(?:\w+)?\n(.*?)```", result, re.DOTALL)
            code_to_save = code_match.group(1) if code_match else result
            out.write_text(code_to_save, encoding="utf-8")
            return f"{result[:400]}...\n\nSaved: {out.name}"
        return result
    except Exception as e:
        return f"AI processing failed: {e}"

def _process_audio(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "transcribe"

    if action == "info":
        try:
            from pydub import AudioSegment
            audio    = AudioSegment.from_file(path)
            duration = len(audio) / 1000
            mins, secs = divmod(int(duration), 60)
            return (f"Audio: {mins}m {secs}s, "
                    f"{audio.channels} ch, "
                    f"{audio.frame_rate}Hz, "
                    f"{_file_size_str(path)}")
        except ImportError:
            return f"Audio file: {_file_size_str(path)} (install pydub for more info)"
        except Exception as e:
            return f"Info failed: {e}"

    if action == "transcribe":
        try:
            model   = _gemini_client()
            content = path.read_bytes()
            mime    = {
                "mp3": "audio/mp3", "wav": "audio/wav",
                "ogg": "audio/ogg", "m4a": "audio/mp4",
                "aac": "audio/aac", "flac": "audio/flac",
            }.get(path.suffix.lstrip(".").lower(), "audio/mpeg")
            response = model.generate_content([
                "Transcribe all speech in this audio file accurately.",
                {"mime_type": mime, "data": content}
            ])
            result = response.text.strip()
            if params.get("save", True):
                out = _output_path(path, "transcript", ".txt")
                out.write_text(result, encoding="utf-8")
                return f"Transcription saved: {out.name}\n\nPreview: {result[:300]}"
            return result
        except Exception as e:
            return f"Transcription failed: {e}"

    if action == "convert":
        fmt = params.get("format", "mp3").lstrip(".")
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(path)
            out   = _output_path(path, "converted", f".{fmt}")
            audio.export(out, format=fmt)
            return f"Converted to {fmt.upper()}. Saved: {out.name}"
        except ImportError:
            return "pydub not installed. Run: pip install pydub"
        except Exception as e:
            return f"Convert failed: {e}"

    if action == "trim":
        start = float(params.get("start", 0))
        end   = float(params.get("end",   0))
        try:
            from pydub import AudioSegment
            audio   = AudioSegment.from_file(path)
            end_ms  = int(end * 1000)   if end   else len(audio)
            trimmed = audio[int(start * 1000):end_ms]
            out     = _output_path(path, f"trim_{int(start)}s_{int(end)}s")
            trimmed.export(out, format=path.suffix.lstrip("."))
            return f"Trimmed audio ({int(start)}s–{int(end)}s). Saved: {out.name}"
        except ImportError:
            return "pydub not installed."
        except Exception as e:
            return f"Trim failed: {e}"

    return f"Unknown audio action: '{action}'. Try: transcribe, info, convert, trim"

def _process_video(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "info"


    def _ffmpeg_available() -> bool:
        try:
            subprocess.run(["ffmpeg", "-version"], capture_output=True, timeout=3)
            return True
        except Exception:
            return False

    if action == "info":
        try:
            result = subprocess.run(
                ["ffprobe", "-v", "quiet", "-print_format", "json",
                 "-show_format", "-show_streams", str(path)],
                capture_output=True, text=True, timeout=10
            )
            data     = json.loads(result.stdout)
            fmt      = data.get("format", {})
            duration = float(fmt.get("duration", 0))
            mins, secs = divmod(int(duration), 60)
            size     = _file_size_str(path)
            streams  = data.get("streams", [])
            video_s  = next((s for s in streams if s["codec_type"] == "video"), {})
            w        = video_s.get("width", "?")
            h        = video_s.get("height", "?")
            fps      = video_s.get("r_frame_rate", "?")
            return f"Video: {mins}m {secs}s, {w}x{h}, {fps} fps, {size}"
        except Exception:
            return f"Video file: {_file_size_str(path)}"

    if action == "extract_audio":
        if not _ffmpeg_available():
            return "ffmpeg not found. Install ffmpeg to extract audio."
        out = _output_path(path, "audio", ".mp3")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), "-q:a", "0", "-map", "a", str(out), "-y"],
                capture_output=True, timeout=300
            )
            return f"Audio extracted. Saved: {out.name}"
        except Exception as e:
            return f"Extract audio failed: {e}"

    if action == "trim":
        start = params.get("start", "00:00:00")
        end   = params.get("end",   "")
        if not _ffmpeg_available():
            return "ffmpeg not found."
        out = _output_path(path, f"trim", path.suffix)
        try:
            cmd = ["ffmpeg", "-i", str(path), "-ss", str(start)]
            if end:
                cmd += ["-to", str(end)]
            cmd += ["-c", "copy", str(out), "-y"]
            subprocess.run(cmd, capture_output=True, timeout=600)
            return f"Trimmed video saved: {out.name}"
        except Exception as e:
            return f"Trim failed: {e}"

    if action == "extract_frame":
        timestamp = params.get("timestamp", "00:00:01")
        if not _ffmpeg_available():
            return "ffmpeg not found."
        out = _output_path(path, f"frame_{timestamp.replace(':', '')}", ".jpg")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), "-ss", timestamp,
                 "-vframes", "1", str(out), "-y"],
                capture_output=True, timeout=30
            )
            return f"Frame extracted at {timestamp}. Saved: {out.name}"
        except Exception as e:
            return f"Extract frame failed: {e}"

    if action == "compress":
        crf = int(params.get("quality", 28))  
        if not _ffmpeg_available():
            return "ffmpeg not found."
        out = _output_path(path, f"compressed_crf{crf}", ".mp4")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path),
                 "-c:v", "libx264", "-crf", str(crf),
                 "-preset", "medium", "-c:a", "copy",
                 str(out), "-y"],
                capture_output=True, timeout=1800
            )
            before = _file_size_str(path)
            after  = _file_size_str(out)
            return f"Compressed: {before} → {after}. Saved: {out.name}"
        except Exception as e:
            return f"Compress failed: {e}"

    if action == "transcribe":
        if not _ffmpeg_available():
            return "ffmpeg not found. Needed for video transcription."
        tmp_audio = Path(tempfile.mktemp(suffix=".mp3"))
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), "-q:a", "0", "-map", "a",
                 str(tmp_audio), "-y"],
                capture_output=True, timeout=300
            )
            result = _process_audio(tmp_audio, "transcribe", params, speak)
            return result
        except Exception as e:
            return f"Video transcription failed: {e}"
        finally:
            if tmp_audio.exists():
                tmp_audio.unlink()

    if action == "convert":
        fmt = params.get("format", "mp4").lstrip(".")
        if not _ffmpeg_available():
            return "ffmpeg not found."
        out = _output_path(path, "converted", f".{fmt}")
        try:
            subprocess.run(
                ["ffmpeg", "-i", str(path), str(out), "-y"],
                capture_output=True, timeout=1800
            )
            return f"Converted to {fmt.upper()}. Saved: {out.name}"
        except Exception as e:
            return f"Convert failed: {e}"

    return f"Unknown video action: '{action}'. Try: info, trim, extract_audio, extract_frame, compress, transcribe, convert"

def _process_archive(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "list"

    if action == "list":
        try:
            import zipfile, tarfile
            ext = path.suffix.lower()
            if ext == ".zip":
                with zipfile.ZipFile(path) as z:
                    names = z.namelist()
            elif ext in (".tar", ".gz", ".bz2", ".xz"):
                with tarfile.open(path) as t:
                    names = t.getnames()
            else:
                return f"Unsupported archive format: {ext}"
            preview = "\n".join(names[:30])
            suffix  = f"\n... and {len(names)-30} more" if len(names) > 30 else ""
            return f"Archive contains {len(names)} files:\n{preview}{suffix}"
        except Exception as e:
            return f"List failed: {e}"

    if action == "extract":
        dest = Path(params.get("destination", str(path.parent / path.stem)))
        dest.mkdir(parents=True, exist_ok=True)
        try:
            shutil.unpack_archive(path, dest)
            return f"Extracted to: {dest}"
        except Exception as e:
            return f"Extract failed: {e}"

    return f"Unknown archive action: '{action}'. Try: list, extract"

def _process_pptx(path: Path, action: str, params: dict, speak=None) -> str:
    action = action or "summarize"
    style = params.get("style") or params.get("tone") or params.get("essay_style") or "student"

    if action in ("summarize", "extract_text", "analyze"):
        text, image_notes = _extract_pptx_text_and_images(path)
        combined = _merge_text_and_image_notes(text, image_notes)
        if action == "extract_text":
            out = _output_path(path, "text", ".txt")
            out.write_text(combined, encoding="utf-8")
            return f"Text and visual notes extracted. Saved: {out.name}"
        try:
            model    = _gemini_client()
            prompt   = _essay_prompt(
                "presentation",
                combined[:30000],
                "Write a polished sample essay about this presentation. Describe the main argument or story, the most important slides, notable image-based content, and the overall takeaway. Use a title, thesis, body paragraphs, and a short conclusion."
                ,
                style,
            )
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"AI processing failed: {e}"

    return f"Unknown PPTX action: '{action}'. Try: summarize, extract_text, analyze"


def _normalize_file_paths(parameters: dict) -> list[Path]:
    raw_paths = parameters.get("file_paths")
    collected: list[str] = []

    if isinstance(raw_paths, str):
        raw_paths = [raw_paths]

    if isinstance(raw_paths, (list, tuple, set)):
        for item in raw_paths:
            if item is None:
                continue
            text = str(item).strip()
            if text:
                collected.append(text)

    single_path = str(parameters.get("file_path", "") or "").strip()
    if single_path:
        collected.append(single_path)

    paths: list[Path] = []
    seen: set[str] = set()
    for raw in collected:
        path = Path(raw).expanduser()
        key = str(path.resolve()) if path.exists() else str(path)
        if key in seen:
            continue
        seen.add(key)
        if path.exists() and path.is_file():
            paths.append(path)
    return paths


def _bundle_action_for_type(file_type: str) -> str:
    return {
        "image": "analyze",
        "pdf": "summarize",
        "docx": "summarize",
        "text": "summarize",
        "csv": "analyze",
        "excel": "analyze",
        "json": "analyze",
        "xml": "analyze",
        "code": "summarize",
        "audio": "transcribe",
        "video": "transcribe",
        "archive": "list",
        "pptx": "summarize",
    }.get(file_type, "summarize")


def _summarize_file_for_bundle(path: Path, action: str, params: dict) -> str:
    file_type = _detect_type(path)
    effective_action = (action or _bundle_action_for_type(file_type) or "summarize").lower().strip()
    shared_params = {**params, "save": False}

    try:
        if file_type == "image":
            try:
                from PIL import Image

                model = _gemini_client()
                img = Image.open(path)
                prompt = (
                    "Create a concise comparison-ready summary of this image. "
                    "Mention the subject, visible text, notable objects, colours, layout, and any obvious theme. "
                    "Keep it structured and brief."
                )
                if params.get("instruction"):
                    prompt = str(params["instruction"])
                response = model.generate_content([prompt, img])
                return response.text.strip()
            except Exception as e:
                return f"Image summary failed: {e}"

        if file_type == "pdf":
            return _process_pdf(path, "summarize" if effective_action in {"analyze", "summarize", "compare", "cross_analyze"} else effective_action, shared_params)
        if file_type == "docx":
            return _process_text_doc(path, "docx", "summarize" if effective_action in {"analyze", "summarize", "compare", "cross_analyze"} else effective_action, shared_params)
        if file_type == "text":
            return _process_text_doc(path, "text", "summarize" if effective_action in {"analyze", "summarize", "compare", "cross_analyze"} else effective_action, shared_params)
        if file_type == "csv":
            return _process_data(path, "csv", "analyze" if effective_action in {"analyze", "summarize", "compare", "cross_analyze"} else effective_action, shared_params)
        if file_type == "excel":
            return _process_data(path, "excel", "analyze" if effective_action in {"analyze", "summarize", "compare", "cross_analyze"} else effective_action, shared_params)
        if file_type == "json":
            return _process_json(path, "analyze" if effective_action in {"analyze", "summarize", "compare", "cross_analyze"} else effective_action, shared_params)
        if file_type == "xml":
            return _process_json(path, "analyze" if effective_action in {"analyze", "summarize", "compare", "cross_analyze"} else effective_action, shared_params)
        if file_type == "code":
            return _process_code(path, "summarize" if effective_action in {"analyze", "summarize", "compare", "cross_analyze"} else effective_action, shared_params)
        if file_type == "audio":
            chosen = "transcribe" if effective_action in {"analyze", "summarize", "compare", "cross_analyze", "transcribe"} else "info"
            return _process_audio(path, chosen, shared_params)
        if file_type == "video":
            chosen = "transcribe" if effective_action in {"analyze", "summarize", "compare", "cross_analyze", "transcribe"} else "info"
            return _process_video(path, chosen, shared_params)
        if file_type == "archive":
            chosen = "list" if effective_action in {"analyze", "summarize", "compare", "cross_analyze", "list"} else effective_action
            return _process_archive(path, chosen, shared_params)
        if file_type == "pptx":
            return _process_pptx(path, "summarize" if effective_action in {"analyze", "summarize", "compare", "cross_analyze"} else effective_action, shared_params)

        preview = path.read_text(encoding="utf-8", errors="ignore")[:10000]
        try:
            model = _gemini_client()
            prompt = (
                f"File: {path.name} ({path.suffix or 'unknown'})\n"
                f"Content preview:\n{preview}\n\n"
                "Create a concise comparison-ready summary for this uploaded file. "
                "Mention the main topic, recurring ideas, standout values, and any useful context for cross-file comparison."
            )
            if params.get("instruction"):
                prompt = f"{params['instruction']}\n\n{prompt}"
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"Unknown file type ({path.suffix}). Could not process: {e}"
    except Exception as e:
        return f"Bundle summary failed for {path.name}: {e}"


def _cross_analyze_bundle(paths: list[Path], action: str, params: dict, speak=None) -> str:
    style = params.get("style") or params.get("tone") or params.get("essay_style") or "student"
    sections: list[str] = []
    for index, path in enumerate(paths, 1):
        file_type = _detect_type(path)
        summary = _summarize_file_for_bundle(path, action, params)
        sections.append(
            f"[{index}] {path.name} | type={file_type} | size={_file_size_str(path)}\n{summary[:4000]}"
        )

    bundle_prompt = (
        f"You are comparing multiple uploaded files. Write a polished comparative essay in a {style} style based on the summaries below. "
        "Explain the shared themes, common patterns, contradictions, differences, and any outliers in clear prose. "
        "Use a short title, an opening thesis, 2 to 4 body paragraphs, and a brief concluding paragraph. "
        "If helpful, you may finish with a short 'File-by-File Notes' section, but keep the main response essay-like and readable.\n\n"
        + "\n\n".join(sections)
    )

    instruction = str(params.get("instruction", "") or "").strip()
    if instruction:
        bundle_prompt = f"{instruction}\n\n{bundle_prompt}"

    model = _gemini_client()
    response = model.generate_content(bundle_prompt)
    return response.text.strip()

def file_processor(parameters: dict, player=None, speak=None) -> str:
    action      = (parameters.get("action") or "").lower().strip()
    instruction = parameters.get("instruction", "")
    params      = {**parameters, "instruction": instruction}

    paths = _normalize_file_paths(parameters)
    if not paths:
        return "No file path provided."

    if len(paths) > 1:
        log_msg = f"[FileProcessor] MULTI | {len(paths)} files | action={action or 'auto'}"
        print(log_msg)
        if player:
            player.write_log(log_msg)
        try:
            return _cross_analyze_bundle(paths, action or "cross_analyze", params, speak) or "Done."
        except Exception as e:
            import traceback
            traceback.print_exc()
            return f"Processing failed: {e}"

    path = paths[0]
    file_type = _detect_type(path)

    log_msg = f"[FileProcessor] {file_type.upper()} | {path.name} | action={action or 'auto'}"
    print(log_msg)
    if player:
        player.write_log(log_msg)

    if file_type == "unknown":
        try:
            content = path.read_text(encoding="utf-8", errors="ignore")[:10000]
            model   = _gemini_client()
            prompt  = f"File: {path.name}\nContent preview:\n{content}\n\nTask: {action or instruction or 'Describe what this file contains and what can be done with it.'}"
            response = model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return f"Unknown file type ({path.suffix}). Could not process: {e}"

    dispatch = {
        "image":   _process_image,
        "pdf":     _process_pdf,
        "docx":    lambda p, a, pm, s: _process_text_doc(p, "docx", a, pm, s),
        "text":    lambda p, a, pm, s: _process_text_doc(p, "text", a, pm, s),
        "csv":     lambda p, a, pm, s: _process_data(p, "csv",   a, pm, s),
        "excel":   lambda p, a, pm, s: _process_data(p, "excel", a, pm, s),
        "json":    _process_json,
        "xml":     lambda p, a, pm, s: _process_json(p, a, pm, s),  
        "code":    _process_code,
        "audio":   _process_audio,
        "video":   _process_video,
        "archive": _process_archive,
        "pptx":    _process_pptx,
    }

    handler = dispatch.get(file_type)
    if not handler:
        return f"Unsupported file type: {file_type}"

    try:
        result = handler(path, action, params, speak)
        return result or "Done."
    except Exception as e:
        import traceback
        traceback.print_exc()
        return f"Processing failed: {e}"