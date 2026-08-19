import json
import os
import re
from datetime import datetime, timezone
from pathlib import Path

from memory.conversation_log import _cosine, embed_text


BASE_DIR = Path(__file__).resolve().parent.parent
DATA_DIR = Path(os.getenv("JARVIS_DATA_DIR") or (BASE_DIR / "memory")).expanduser()
INDEX_PATH = DATA_DIR / "document_index.json"

_TEXT_EXTENSIONS = {
    ".md", ".markdown", ".txt", ".rst", ".log",
    ".py", ".pyw", ".js", ".jsx", ".ts", ".tsx",
    ".html", ".css", ".scss", ".json", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".ini", ".csv", ".tsv",
    ".sql", ".java", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".go", ".rs", ".rb", ".php", ".swift",
    ".kt", ".kts", ".sh", ".bash", ".ps1", ".lua",
    ".r", ".m", ".mm", ".xml", ".svg"
}

_POWERPOINT_EXTENSIONS = {".ppt", ".pptx"}


def _safe_read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""


def _normalize_text(raw: str) -> str:
    text = raw.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _tokenize(text: str) -> set[str]:
    cleaned = re.sub(r"[^a-z0-9_\- ]+", " ", text.lower())
    return {token.strip() for token in cleaned.split() if token.strip()}


def _chunk_text(text: str, chunk_size: int = 1200, overlap: int = 160) -> list[str]:
    if not text:
        return []
    normalized = text.strip()
    if len(normalized) <= chunk_size:
        return [normalized]

    chunks: list[str] = []
    start = 0
    while start < len(normalized):
        end = min(len(normalized), start + chunk_size)
        chunk = normalized[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(normalized):
            break
        start = max(start + chunk_size - overlap, end - overlap)
    return chunks


def _load_index() -> dict:
    if not INDEX_PATH.exists():
        return {"documents": []}
    try:
        data = json.loads(INDEX_PATH.read_text(encoding="utf-8"))
        if isinstance(data, dict) and isinstance(data.get("documents"), list):
            return data
    except Exception:
        pass
    return {"documents": []}


def _save_index(index: dict) -> None:
    INDEX_PATH.parent.mkdir(parents=True, exist_ok=True)
    INDEX_PATH.write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")


def _extract_pdf_text(path: Path) -> str:
    try:
        import pymupdf
    except Exception:
        return ""

    try:
        doc = pymupdf.open(path)
        pages: list[str] = []
        for page in doc:
            page_text = page.get_text("text") or ""
            if page_text.strip():
                pages.append(page_text.strip())
        doc.close()
        return "\n\n".join(pages)
    except Exception:
        return ""


def _extract_ppt_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return ""

    try:
        presentation = Presentation(path)
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides, 1):
            parts = [f"--- Slide {index} ---"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    parts.append(shape.text.strip())
            slide_text = "\n".join(parts)
            if slide_text.strip():
                slides.append(slide_text)
        return "\n\n".join(slides)
    except Exception:
        return _safe_read_text(path)


def _extract_document_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf_text(path)
    if suffix in {".docx", ".doc"}:
        try:
            from docx import Document
            doc = Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n".join(paragraphs)
        except Exception:
            return _safe_read_text(path)
    if suffix in _POWERPOINT_EXTENSIONS:
        return _extract_ppt_text(path)
    return _safe_read_text(path)


def _index_document_entry(path: Path, source_name: str | None = None) -> dict:
    text = _extract_document_text(path)
    normalized = _normalize_text(text)
    if not normalized:
        return {
            "status": "empty",
            "path": str(path),
            "source_name": source_name or path.name,
            "chunks": 0,
        }

    chunks = []
    for idx, chunk in enumerate(_chunk_text(normalized), 1):
        chunks.append({
            "id": f"{path.name}:{idx}",
            "content": chunk,
        })

    return {
        "status": "indexed",
        "path": str(path),
        "source_name": source_name or path.stem,
        "kind": path.suffix.lower() or "file",
        "added_at": datetime.now(timezone.utc).isoformat(timespec="seconds").replace("+00:00", "Z"),
        "chunks": chunks,
    }


def ingest_document(file_path: str, source_name: str | None = None) -> dict:
    path = Path(file_path).expanduser().resolve()
    if not path.exists():
        return {"status": "missing", "path": str(path), "error": "File not found."}

    if path.is_dir():
        return index_codebase(str(path), source_name=source_name)

    index = _load_index()
    entry = _index_document_entry(path, source_name)
    if entry["status"] == "empty":
        return entry

    documents = [doc for doc in index.get("documents", []) if doc.get("path") != str(path)]
    documents.append(entry)
    index["documents"] = documents
    _save_index(index)
    return {"status": "indexed", "path": str(path), "source_name": entry["source_name"], "chunks": len(entry["chunks"])}


def index_codebase(root_path: str, source_name: str | None = None) -> dict:
    root = Path(root_path).expanduser().resolve()
    if not root.exists():
        return {"status": "missing", "path": str(root), "error": "Directory not found."}

    candidates = []
    for file_path in sorted(root.rglob("*")):
        if file_path.is_dir():
            continue
        suffix = file_path.suffix.lower()
        if suffix in _TEXT_EXTENSIONS or suffix == ".pdf" or suffix in _POWERPOINT_EXTENSIONS:
            candidates.append(file_path)

    index = _load_index()
    kept: list[dict] = []
    seen_paths: set[str] = set()
    for document in index.get("documents", []):
        doc_path = str(document.get("path", ""))
        if doc_path and doc_path.startswith(str(root)):
            continue
        kept.append(document)
        seen_paths.add(doc_path)

    for file_path in candidates:
        if str(file_path) in seen_paths:
            continue
        entry = _index_document_entry(file_path, source_name or file_path.stem)
        if entry.get("status") == "indexed":
            kept.append(entry)

    index["documents"] = kept
    _save_index(index)
    return {"status": "indexed", "path": str(root), "count": len(kept), "files": len(candidates)}


def search_document_index(query: str, limit: int = 5) -> list[dict]:
    q = (query or "").strip()
    if not q:
        return []

    tokens = _tokenize(q)
    if not tokens:
        return []

    matches: list[dict] = []
    index = _load_index()
    for document in index.get("documents", []):
        for chunk in document.get("chunks", []):
            content = str(chunk.get("content", "")).strip()
            if not content:
                continue
            score = 0
            for token in tokens:
                if token in content.lower():
                    score += 1
            if score == 0:
                continue
            snippet = content[:500].strip()
            matches.append({
                "source": document.get("source_name") or document.get("path"),
                "path": document.get("path"),
                "score": score,
                "content": snippet,
            })

    matches.sort(key=lambda item: item["score"], reverse=True)
    return matches[:limit]


def recall_document_details(query: str, limit: int = 3) -> str:
    matches = search_document_index(query, limit=limit)
    if not matches:
        return f"I could not find relevant details for '{query}' in the local document index."

    parts = [f"Relevant local document matches for '{query}':"]
    for idx, match in enumerate(matches, 1):
        source = match.get("source") or "document"
        content = match.get("content", "").strip()
        parts.append(f"{idx}. {source}\n{content}")
    return "\n\n".join(parts)
